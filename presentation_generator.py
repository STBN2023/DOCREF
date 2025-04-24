# File: presentation_generator.py

import os
from io import BytesIO
from pptx import Presentation
from config import PLACEHOLDERS_DEFAULT, TYPES_FORMATAGE, MAPPING_BLASONS
from logger_config import get_logger
from utils.data_utils import convertir_en_numerique, formater_valeur
from utils.text_utils import analyze_placeholders, find_placeholders_in_slide, normalize_placeholder, remplacer_placeholder
from utils.image_utils import creer_image_transparente, remplacer_image
from utils.blasons_utils import gerer_blasons_ameliore


def generate_presentation(model_path, df, placeholders_mapping, img_dir, logos_dir, valeur_blason_active):
    """
    Génère une présentation PPTX en remplaçant placeholders, images et blasons,
    en appliquant le format monétaire pour {{Montant_Travaux}}.
    """
    logger = get_logger(__name__)
    prs = Presentation(model_path)
    # Image transparente pour blasons inactifs
    creer_image_transparente()

    # Réinitialiser l'index pour faire correspondre séquentiellement les slides filtrées
    df = df.reset_index(drop=True)

    for idx, projet in df.iterrows():
        if idx >= len(prs.slides):
            break
        slide = prs.slides[idx]
        slide_phs = find_placeholders_in_slide(slide)

        # Remplacement des placeholders texte
        for ph, col in placeholders_mapping.items():
            if not col:
                continue
            norm = normalize_placeholder(ph)
            actual = slide_phs.get(norm) or next((v for k, v in slide_phs.items() if k.lower() == norm.lower()), None)
            if not actual:
                continue

            val = projet.get(col, '')
            val = convertir_en_numerique(val)
            # Utilisation du mapping TYPES_FORMATAGE basé sur la clé originale ph
            type_format = TYPES_FORMATAGE.get(ph)
            fmt = formater_valeur(val, type_format)
            remplacer_placeholder(slide, actual, fmt)

        # Remplacement de l'image projet
        img_col = placeholders_mapping.get('IMAGE_PROJET')
        if img_col:
            num = ''.join(filter(str.isdigit, str(projet.get(img_col, ''))))
            for ext in ['png', 'jpg', 'jpeg']:
                path = os.path.join(img_dir, f"{num}.{ext}")
                if os.path.exists(path):
                    remplacer_image(slide, 'IMAGE_PROJET', path)
                    break

        # Gestion des blasons
        gerer_blasons_ameliore(slide, projet, MAPPING_BLASONS, logos_dir, valeur_blason_active)

    # Sauvegarde
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
